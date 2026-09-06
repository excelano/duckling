#!/usr/bin/env python3
"""Author the demo documents the Store listing screenshots are taken against.

WHY THESE EXIST AT ALL

`packaging/store-listing.md` asks the screenshots to show a queue with a mix of
formats. The obvious source was docling.rs's own test corpus, and David decided
against it on 2026-09-05: those files are real third-party documents - arXiv
papers, somebody's handbook - and putting their content in a commercial store
listing is a licence question nobody wants to answer later. These are ours, they
are invented, and the organisation in them does not exist.

Six documents, chosen so the queue column reads as the kind of work a person
actually converts, and so that between them they exercise what the listing
claims: headings, tables, lists, pictures, and a page with no text layer at all.

WHY THIS IS PYTHON IN A RUST REPOSITORY

Nothing in the build runs it and nothing in CI runs it, which is the same
standing `packaging/windows/make-ico` has - its own thing, off to one side, so
that what it needs never reaches the shipped binary. The difference is the
language, and that is a measurement rather than a preference: this machine has
no Office, no LibreOffice and no pandoc, and Rust has no PowerPoint writer worth
taking a dependency on. Python has a mature library per format. The generated
documents are committed beside this file so that the Windows and macOS lanes
photograph the same six, and so that running this is never a precondition for
anything.

    python -m venv .venv
    .venv/Scripts/pip install -r requirements.txt      # Scripts on Windows, bin elsewhere
    .venv/Scripts/python make-demo-docs.py

Author: David M. Anderson
Built with AI assistance (Claude, Anthropic)
"""

import random
import zipfile
from pathlib import Path

from docx import Document
from docx.shared import Pt as DocxPt
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font as XlFont
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_JUSTIFY
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

HERE = Path(__file__).resolve().parent
OUT = HERE / "documents"

# The invented organisation everything below belongs to.
#
# Deliberately mundane and deliberately not anybody: a small field station doing
# water monitoring. It gives tables of real-looking measurements, which is what
# shows off table extraction, and its subject matter cannot embarrass anyone in
# a screenshot the way invented business figures or an invented invoice could.
STATION = "Alder Creek Field Station"

# The readings every document shares, so that the six look like one body of work
# rather than six unrelated samples. A person who reads two of them in a
# screenshot should find they agree.
READINGS = [
    ("AC-01", "Upper meadow", "12 Mar", "7.1", "9.4", "142"),
    ("AC-02", "Culvert outflow", "12 Mar", "6.8", "8.1", "168"),
    ("AC-03", "Beaver pond", "13 Mar", "7.4", "6.2", "121"),
    ("AC-04", "Willow bend", "13 Mar", "7.2", "8.8", "134"),
    ("AC-05", "Lower gauge", "14 Mar", "6.9", "9.1", "155"),
]
READING_HEAD = ["Site", "Location", "Date", "pH", "DO mg/L", "Cond. uS/cm"]

# The same five sites in November, so the second report is a different document
# rather than the first one with the month changed. The conductivity column is
# where the two are meant to be read against each other.
WINTER_READINGS = [
    ("AC-01", "Upper meadow", "14 Nov", "7.2", "10.1", "138"),
    ("AC-02", "Culvert outflow", "14 Nov", "7.0", "9.6", "121"),
    ("AC-03", "Beaver pond", "15 Nov", "7.3", "7.4", "119"),
    ("AC-04", "Willow bend", "15 Nov", "7.2", "9.9", "130"),
    ("AC-05", "Lower gauge", "16 Nov", "7.1", "10.3", "147"),
]

# The instruments, which give the second Word file a table of its own that is
# not another copy of the readings.
EQUIPMENT = [
    ("Multiparameter probe", "MP-4471", "Mar 2026", "Annual"),
    ("Spare probe", "MP-2210", "Nov 2025", "Annual"),
    ("Buffer standards pH 4", "-", "Jan 2026", "On opening"),
    ("Buffer standards pH 7", "-", "Jan 2026", "On opening"),
    ("Gauge board camera", "GB-08", "Aug 2025", "Two-yearly"),
    ("Field radio", "FR-13", "Feb 2026", "Annual"),
]
EQUIPMENT_HEAD = ["Instrument", "Asset", "Last checked", "Interval"]

SUMMARY = (
    "Five monitoring sites were visited over three days in March. Dissolved "
    "oxygen remained above the 6.0 mg/L threshold at every site, and pH stayed "
    "within the range the station has recorded since monitoring began. "
    "Conductivity at the culvert outflow is the one reading worth returning to: "
    "it has risen in each of the last four surveys and is now the highest on "
    "the reach."
)

METHOD = (
    "Readings are taken mid-channel at two thirds depth, upstream of the "
    "observer, with the probe allowed to settle for sixty seconds before the "
    "value is recorded. Where flow is too shallow for the probe, a grab sample "
    "is taken and measured on the bank within ten minutes. Every instrument is "
    "checked against buffer standards on the morning of a survey and again on "
    "return."
)


def report_pdf(path, title="Spring Water Quality Survey", period="March 2026",
               readings=None, summary=None, closing=None):
    """A multi-page PDF with a real text layer, headings and a table.

    The one that is worth photographing mid-conversion: it is long enough that
    the queue shows a page count while the layout model works through it.

    Parameterised so the winter survey is the same kind of document rather than
    a second function that would drift from this one.
    """
    readings = readings if readings is not None else READINGS
    summary = summary if summary is not None else SUMMARY
    closing = closing if closing is not None else (
        "Conductivity at AC-02 has risen from 121 to 168 microsiemens over "
        "four surveys. The culvert drains the lane and the verge above it, "
        "and the rise is consistent with road salt still working through "
        "after the winter. The reading is not itself a concern at this "
        "level; the trend is what the next survey should settle."
    )
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle(
        "h1", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=16,
        spaceAfter=12, textColor=colors.HexColor("#2f3337"),
    )
    h2 = ParagraphStyle(
        "h2", parent=styles["Heading2"], fontName="Helvetica-Bold", fontSize=12,
        spaceBefore=14, spaceAfter=6, textColor=colors.HexColor("#4a6fa5"),
    )
    body = ParagraphStyle(
        "body", parent=styles["BodyText"], fontSize=10.5, leading=15,
        alignment=TA_JUSTIFY, spaceAfter=8,
    )

    story = [
        Paragraph(title, h1),
        Paragraph(f"{STATION} &middot; Reach 4 &middot; {period}", body),
        Paragraph("Summary", h2),
        Paragraph(summary, body),
        Paragraph("Method", h2),
        Paragraph(METHOD, body),
        Paragraph("Readings", h2),
    ]

    table = Table(
        [READING_HEAD] + [list(r) for r in readings],
        colWidths=[0.8 * inch, 1.9 * inch, 0.9 * inch, 0.7 * inch, 1.1 * inch, 1.3 * inch],
    )
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#4a6fa5")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 9.5),
        ("ALIGN", (3, 1), (-1, -1), "RIGHT"),
        ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#b8c4d4")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f2f5f9")]),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
    ]))
    story += [table, Spacer(1, 14)]

    story += [
        Paragraph("Notes on the culvert outflow", h2),
        Paragraph(closing, body),
        PageBreak(),
        Paragraph("Observations by site", h1),
    ]
    for site, place, date, ph, do, cond in readings:
        story += [
            Paragraph(f"{site} &mdash; {place}", h2),
            Paragraph(
                f"Visited {date}. pH {ph}, dissolved oxygen {do} mg/L, "
                f"conductivity {cond} microsiemens per centimetre. Channel "
                "clear, no odour, no visible film. Bank vegetation intact and "
                "the gauge board legible from the near bank.",
                body,
            ),
        ]
    story += [
        Paragraph("Recommendations", h2),
        Paragraph(
            "Continue the four-weekly cycle through the summer. Add a second "
            "reading above the culvert so the rise at AC-02 can be attributed "
            "rather than inferred. Replace the gauge board at AC-05, which is "
            "legible but weathered.",
            body,
        ),
    ]

    SimpleDocTemplate(
        str(path), pagesize=LETTER,
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.9 * inch, bottomMargin=0.9 * inch,
        title=title, author=STATION,
    ).build(story)


def notes_docx(path):
    """A Word file with headings, a bulleted list and a table."""
    doc = Document()
    doc.core_properties.title = "Field Notes"
    doc.core_properties.author = STATION

    doc.add_heading("Field Notes", level=0)
    doc.add_paragraph(f"{STATION} \u00b7 Reach 4 \u00b7 12-14 March 2026")

    doc.add_heading("Conditions", level=1)
    doc.add_paragraph(
        "Cold and bright for all three days. No rain in the seventy-two hours "
        "before the first visit, so flows were low and the readings should be "
        "read as a dry-weather baseline."
    )

    doc.add_heading("Carried", level=1)
    for item in [
        "Multiparameter probe and spare batteries",
        "Buffer standards, pH 4 and pH 7",
        "Grab bottles, twelve, labelled",
        "Gauge board camera and the site key",
        "This notebook, and a pencil that works when wet",
    ]:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_heading("Readings", level=1)
    table = doc.add_table(rows=1, cols=len(READING_HEAD))
    table.style = "Light Grid Accent 1"
    for cell, head in zip(table.rows[0].cells, READING_HEAD):
        cell.text = head
        for run in cell.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = DocxPt(9)
    for row in READINGS:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            cell.text = value
            for run in cell.paragraphs[0].runs:
                run.font.size = DocxPt(9)

    doc.add_heading("To follow up", level=1)
    doc.add_paragraph(
        "AC-02 conductivity again. Fourth rise running. Ask about the gritting "
        "schedule for the lane before drawing any conclusion from it."
    )
    doc.add_paragraph(
        "Gauge board at AC-05 is weathered. Readable this time. Order a "
        "replacement before it is not."
    )
    doc.save(str(path))


def log_xlsx(path):
    """A workbook with two sheets, because a spreadsheet backend that only ever
    sees one sheet is not being shown doing anything."""
    book = Workbook()

    sheet = book.active
    sheet.title = "Readings"
    sheet["A1"] = "Spring Water Quality Survey"
    sheet["A1"].font = XlFont(bold=True, size=14)
    sheet["A2"] = f"{STATION} \u00b7 Reach 4 \u00b7 March 2026"
    sheet.append([])
    sheet.append(READING_HEAD)
    for cell in sheet[4]:
        cell.font = XlFont(bold=True)
        cell.alignment = Alignment(horizontal="center")
    for site, place, date, ph, do, cond in READINGS:
        sheet.append([site, place, date, float(ph), float(do), int(cond)])
    for column, width in zip("ABCDEF", (10, 22, 12, 8, 12, 14)):
        sheet.column_dimensions[column].width = width

    trend = book.create_sheet("Conductivity trend")
    trend["A1"] = "AC-02 culvert outflow, microsiemens per centimetre"
    trend["A1"].font = XlFont(bold=True, size=12)
    trend.append([])
    trend.append(["Survey", "Date", "Conductivity"])
    for cell in trend[3]:
        cell.font = XlFont(bold=True)
    for survey, date, value in [
        ("2025-11", "14 Nov 2025", 121),
        ("2025-12", "12 Dec 2025", 134),
        ("2026-01", "16 Jan 2026", 149),
        ("2026-03", "12 Mar 2026", 168),
    ]:
        trend.append([survey, date, value])
    for column, width in zip("ABC", (12, 16, 14)):
        trend.column_dimensions[column].width = width

    book.save(str(path))


def deck_pptx(path):
    """A short deck: title, bullets, and a table slide."""
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)

    title = deck.slides.add_slide(deck.slide_layouts[0])
    title.shapes.title.text = "Spring Water Quality Survey"
    title.placeholders[1].text = f"{STATION} \u00b7 Reach 4 \u00b7 March 2026"

    findings = deck.slides.add_slide(deck.slide_layouts[1])
    findings.shapes.title.text = "What the survey found"
    body = findings.placeholders[1].text_frame
    body.text = "Dissolved oxygen above threshold at all five sites"
    for line in [
        "pH within the range recorded since monitoring began",
        "Conductivity at the culvert outflow up for a fourth survey",
        "No visible film, odour or bank damage at any site",
    ]:
        paragraph = body.add_paragraph()
        paragraph.text = line
        paragraph.level = 0

    table_slide = deck.slides.add_slide(deck.slide_layouts[5])
    table_slide.shapes.title.text = "Readings"
    rows, cols = len(READINGS) + 1, len(READING_HEAD)
    shape = table_slide.shapes.add_table(
        rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.4 * rows)
    )
    grid = shape.table
    for column, head in enumerate(READING_HEAD):
        cell = grid.cell(0, column)
        cell.text = head
        cell.text_frame.paragraphs[0].runs[0].font.size = PptxPt(14)
    for r, row in enumerate(READINGS, start=1):
        for c, value in enumerate(row):
            cell = grid.cell(r, c)
            cell.text = value
            cell.text_frame.paragraphs[0].runs[0].font.size = PptxPt(13)

    next_slide = deck.slides.add_slide(deck.slide_layouts[1])
    next_slide.shapes.title.text = "Next"
    frame = next_slide.placeholders[1].text_frame
    frame.text = "Keep the four-weekly cycle through the summer"
    for line in [
        "Add a reading above the culvert to attribute the AC-02 rise",
        "Replace the weathered gauge board at AC-05",
    ]:
        paragraph = frame.add_paragraph()
        paragraph.text = line

    deck.save(str(path))


def handbook_epub(path):
    """An EPUB, assembled by hand.

    EPUB is a zip with three fixed parts and some XHTML, so a library for it
    would be a dependency to read rather than to trust. The one rule that bites
    is the first entry: `mimetype`, stored rather than deflated, written first.
    """
    chapters = [
        ("Before you go out", [
            "Sign the board in the porch with your name, your reach and the "
            "hour you expect to be back. If you are not back within two hours "
            "of that, somebody comes looking, and they will start at the last "
            "site on your list.",
            "Take the radio even for a short round. The reach has no coverage "
            "below the culvert and that is exactly where the footing is worst.",
        ]),
        ("Taking a reading", [
            METHOD,
            "If a value looks wrong, take it again before you write it down. If "
            "it is still wrong, write both down and say so in your notes. A "
            "reading nobody can explain is worth more than a reading quietly "
            "dropped.",
        ]),
        ("Coming back", [
            "Rinse the probe in deionised water and cap it wet. A probe stored "
            "dry is a probe that needs recalibrating before it can be trusted, "
            "and that is an hour somebody else will spend.",
            "Enter the readings the same day. The notebook is the record; the "
            "spreadsheet is what everyone else reads.",
        ]),
    ]

    def xhtml(title, paragraphs):
        body = "\n".join(f"    <p>{p}</p>" for p in paragraphs)
        return (
            '<?xml version="1.0" encoding="utf-8"?>\n'
            '<!DOCTYPE html>\n'
            '<html xmlns="http://www.w3.org/1999/xhtml">\n'
            f"  <head><title>{title}</title></head>\n"
            f"  <body>\n    <h1>{title}</h1>\n{body}\n  </body>\n</html>\n"
        )

    items = "\n".join(
        f'    <item id="c{i}" href="c{i}.xhtml" media-type="application/xhtml+xml"/>'
        for i in range(len(chapters))
    )
    spine = "\n".join(f'    <itemref idref="c{i}"/>' for i in range(len(chapters)))
    opf = (
        '<?xml version="1.0" encoding="utf-8"?>\n'
        '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" unique-identifier="id">\n'
        '  <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">\n'
        '    <dc:identifier id="id">urn:uuid:alder-creek-handbook</dc:identifier>\n'
        "    <dc:title>Observer's Handbook</dc:title>\n"
        "    <dc:language>en-GB</dc:language>\n"
        f"    <dc:creator>{STATION}</dc:creator>\n"
        '  </metadata>\n'
        f"  <manifest>\n{items}\n"
        '    <item id="nav" href="nav.xhtml" media-type="application/xhtml+xml" properties="nav"/>\n'
        "  </manifest>\n"
        f"  <spine>\n{spine}\n  </spine>\n"
        "</package>\n"
    )
    links = "\n".join(
        f'      <li><a href="c{i}.xhtml">{t}</a></li>'
        for i, (t, _) in enumerate(chapters)
    )
    nav = (
        '<?xml version="1.0" encoding="utf-8"?>\n'
        '<!DOCTYPE html>\n'
        '<html xmlns="http://www.w3.org/1999/xhtml" xmlns:epub="http://www.idpf.org/2007/ops">\n'
        "  <head><title>Contents</title></head>\n"
        '  <body>\n    <nav epub:type="toc"><h1>Contents</h1>\n      <ol>\n'
        f"{links}\n      </ol>\n    </nav>\n  </body>\n</html>\n"
    )

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as book:
        book.writestr(
            zipfile.ZipInfo("mimetype"), "application/epub+zip",
            compress_type=zipfile.ZIP_STORED,
        )
        book.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0" encoding="utf-8"?>\n'
            '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">\n'
            '  <rootfiles>\n'
            '    <rootfile full-path="OEBPS/content.opf" media-type="application/oebps-package+xml"/>\n'
            "  </rootfiles>\n</container>\n",
        )
        book.writestr("OEBPS/content.opf", opf)
        book.writestr("OEBPS/nav.xhtml", nav)
        for i, (title, paragraphs) in enumerate(chapters):
            book.writestr(f"OEBPS/c{i}.xhtml", xhtml(title, paragraphs))


def scanned_pdf(path):
    """A PDF with no text layer at all: one page, one bitmap, faintly crooked.

    This is the document that proves the models are in the package, and it is
    the only one here that does. Every other file converts through a pure-Rust
    backend and would convert on a machine with no `.models/` beside it; this
    one comes back as text or does not come back at all.

    Made to be read rather than to look convincingly grubby. The tilt and the
    speckle are enough that the page is unmistakably a scan, and light enough
    that the OCR recogniser this application ships still has an easy time -
    a demo that shows OCR failing is not a demo.
    """
    lines = [
        "ALDER CREEK FIELD STATION",
        "",
        "Notice to observers",
        "",
        "The gauge board at the lower site is weathered and a",
        "replacement has been ordered. Until it arrives, record",
        "the staff reading from the upstream face only, and note",
        "the time alongside it.",
        "",
        "The culvert path is closed while the verge is cut. Reach",
        "the outflow from the meadow gate instead. Allow ten",
        "minutes more for the round than the board suggests.",
        "",
        "Posted 9 March 2026",
    ]

    # A real typeface at scanning resolution. The bitmap font PIL falls back on
    # is far too small to be recognised once the page is scaled to letter size,
    # so a missing font here is a failure worth stopping on rather than a
    # quietly worse demo.
    font = None
    for candidate in ("georgia.ttf", "times.ttf", "calibri.ttf", "arial.ttf"):
        try:
            font = ImageFont.truetype(f"C:/Windows/Fonts/{candidate}", 46)
            break
        except OSError:
            continue
    if font is None:
        raise SystemExit("no usable TrueType font found under C:/Windows/Fonts")

    width, height = 1700, 2200            # letter at 200 dpi
    page = Image.new("L", (width, height), 255)
    draw = ImageDraw.Draw(page)
    y = 300
    for line in lines:
        if line:
            draw.text((260, y), line, font=font, fill=28)
        y += 72

    # A scan is never square to the platen and never perfectly white. Both are
    # small on purpose; see the docstring.
    page = page.rotate(-0.6, resample=Image.BICUBIC, fillcolor=255)
    random.seed(20260309)
    pixels = page.load()
    for _ in range(24000):
        x, y = random.randrange(width), random.randrange(height)
        pixels[x, y] = max(0, pixels[x, y] - random.randrange(30, 90))

    bitmap = OUT / "_scan.png"
    page.save(bitmap, "PNG")
    try:
        sheet = pdfcanvas.Canvas(str(path), pagesize=LETTER)
        sheet.drawImage(str(bitmap), 0, 0, width=LETTER[0], height=LETTER[1])
        sheet.showPage()
        sheet.save()
    finally:
        bitmap.unlink(missing_ok=True)


def winter_report_pdf(path):
    """The November survey, so the queue holds two documents of one kind that
    are genuinely different documents."""
    report_pdf(
        path,
        title="Winter Water Quality Survey",
        period="November 2025",
        readings=WINTER_READINGS,
        summary=(
            "Five monitoring sites were visited over three days in November. "
            "Dissolved oxygen was high at every site, as it usually is at this "
            "temperature, and pH was steady across the reach. Conductivity was "
            "unremarkable and is recorded here as the baseline the spring "
            "survey should be read against."
        ),
        closing=(
            "Conductivity at AC-02 was 121 microsiemens, which is the lowest "
            "of the five and in line with the two surveys before it. Nothing "
            "in this survey suggested the site needed watching; that came "
            "later."
        ),
    )


def register_docx(path):
    """A second Word file, with a table that is not another copy of the
    readings: the instruments and when each was last checked."""
    doc = Document()
    doc.core_properties.title = "Equipment Register"
    doc.core_properties.author = STATION

    doc.add_heading("Equipment Register", level=0)
    doc.add_paragraph(f"{STATION} · reviewed March 2026")
    doc.add_paragraph(
        "Every instrument on this list is checked against its interval before "
        "a survey and signed for on return. An instrument past its interval "
        "does not go out."
    )

    table = doc.add_table(rows=1, cols=len(EQUIPMENT_HEAD))
    table.style = "Light Grid Accent 1"
    for cell, head in zip(table.rows[0].cells, EQUIPMENT_HEAD):
        cell.text = head
        for run in cell.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = DocxPt(9)
    for row in EQUIPMENT:
        for cell, value in zip(table.add_row().cells, row):
            cell.text = value
            for run in cell.paragraphs[0].runs:
                run.font.size = DocxPt(9)

    doc.add_heading("Ordered", level=1)
    doc.add_paragraph(
        "A replacement gauge board for AC-05. Quoted six weeks. The existing "
        "board stays readable in the meantime and is photographed on every "
        "visit."
    )
    doc.save(str(path))


def trend_csv(path):
    """A CSV. Plain text, so it exercises the backend that has the least to do
    and the queue row that a person is least likely to expect to work."""
    lines = ["Survey,Date,Site,Conductivity uS/cm,Note"]
    for survey, date, value, note in [
        ("2025-11", "14 Nov 2025", 121, "baseline"),
        ("2025-12", "12 Dec 2025", 134, ""),
        ("2026-01", "16 Jan 2026", 149, "first gritting of the winter"),
        ("2026-03", "12 Mar 2026", 168, "fourth rise running"),
    ]:
        lines.append(f"{survey},{date},AC-02,{value},{note}")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8", newline="\r\n")


def notice_rtf(path):
    """An RTF, written by hand.

    RTF is one of the formats `packaging/store-listing.md` singles out as one
    nothing else touches, so it earns a row. The format is plain text with
    braces; a library for it would be a dependency to read rather than to trust,
    the same call the EPUB above makes.
    """
    paragraphs = [
        r"{\b Notice to observers\par}",
        r"The culvert path is closed while the verge is cut. Reach the "
        r"outflow from the meadow gate instead, and allow ten minutes more "
        r"for the round than the board suggests.\par",
        r"The gauge board at AC-05 is weathered and a replacement has been "
        r"ordered. Until it arrives, record the staff reading from the "
        r"upstream face only and note the time alongside it.\par",
        r"\par",
        r"{\i Posted 9 March 2026}\par",
    ]
    rtf = (
        r"{\rtf1\ansi\deff0"
        r"{\fonttbl{\f0\froman Times New Roman;}}"
        r"\fs24 "
        + "".join(paragraphs)
        + "}"
    )
    path.write_text(rtf, encoding="ascii", newline="")


def index_html(path):
    """HTML, which docling reads structurally rather than as text - so this one
    is here to show headings and a list surviving a format nobody thinks of as a
    document."""
    rows = "\n".join(
        f"      <li><strong>{site}</strong> &mdash; {place}, "
        f"last read {date}</li>"
        for site, place, date, _, _, _ in READINGS
    )
    path.write_text(
        "<!DOCTYPE html>\n"
        '<html lang="en-GB">\n'
        "  <head>\n"
        '    <meta charset="utf-8"/>\n'
        "    <title>Reach 4 site index</title>\n"
        "  </head>\n"
        "  <body>\n"
        "    <h1>Reach 4 site index</h1>\n"
        f"    <p>{STATION}. Five monitoring sites, upstream to downstream.</p>\n"
        "    <h2>Sites</h2>\n"
        f"    <ul>\n{rows}\n    </ul>\n"
        "    <h2>Access</h2>\n"
        "    <p>The meadow gate is the only vehicle access. The lane above the\n"
        "    culvert is not a right of way and the station does not have a key\n"
        "    to it.</p>\n"
        "  </body>\n"
        "</html>\n",
        encoding="utf-8",
    )


def species_md(path):
    """Markdown in, which is worth a row of its own: a person who already has
    Markdown is the person most likely to wonder what this application is for,
    and the answer is that it reads forty other things."""
    path.write_text(
        "# Reach 4 species list\n\n"
        f"{STATION}. Recorded on the March survey, not a complete list.\n\n"
        "## Bankside\n\n"
        "- Alder, *Alnus glutinosa*, dominant along the whole reach\n"
        "- Willow, *Salix* spp., pollarded at the bend\n"
        "- Meadowsweet, in the wet margin above the culvert\n\n"
        "## In the water\n\n"
        "| Group | Sites | Note |\n"
        "| --- | --- | --- |\n"
        "| Mayfly nymph | AC-01, AC-03, AC-04 | absent below the culvert |\n"
        "| Caddis larva | all five | cased, mostly on the riffles |\n"
        "| Freshwater shrimp | AC-02, AC-05 | abundant |\n\n"
        "The absence of mayfly nymph at AC-02 and AC-05 is noted rather than\n"
        "explained. Two surveys is not a trend.\n",
        encoding="utf-8",
    )


# Twelve, and the count is deliberate: six left the queue panel three quarters
# empty in a 1366x768 screenshot, which is the size the Store asks for. Measured
# 2026-09-05. The set is also nine distinct extensions rather than six, which is
# the breadth `store-listing.md` claims doing some work instead of being
# asserted.
DOCUMENTS = [
    ("site-survey-report.pdf", report_pdf),
    ("winter-survey-report.pdf", winter_report_pdf),
    ("field-notes.docx", notes_docx),
    ("equipment-register.docx", register_docx),
    ("sample-log.xlsx", log_xlsx),
    ("orientation-deck.pptx", deck_pptx),
    ("observers-handbook.epub", handbook_epub),
    ("site-index.html", index_html),
    ("species-list.md", species_md),
    ("culvert-trend.csv", trend_csv),
    ("permit-notice.rtf", notice_rtf),
    ("scanned-notice.pdf", scanned_pdf),
]


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    for name, build in DOCUMENTS:
        path = OUT / name
        build(path)
        print(f"wrote {path.relative_to(HERE.parent.parent)}  {path.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
